"""Document intelligence parser router with block-level provenance."""

from __future__ import annotations

import hashlib
import re
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable, List

from ..config import settings
from ..models.schemas import Chunk
from ..parsing.text_quality import normalize_dirty_scientific_text
from .document_models import (
    DocumentIntelligenceResult,
    ParsedBlock,
    ParsedImageRef,
    ParsedTable,
    ScannedPdfDetection,
)
from .parser_audit import ParserAuditWriter

try:
    from docling.document_converter import DocumentConverter
except Exception:  # pragma: no cover
    DocumentConverter = None  # type: ignore
try:
    from markitdown import MarkItDown
except Exception:  # pragma: no cover
    MarkItDown = None  # type: ignore
try:
    import ocrmypdf  # noqa: F401
except Exception:  # pragma: no cover
    ocrmypdf = None  # type: ignore


PARSER_VERSION = "document_intelligence_v1"


@dataclass
class ParsedDocument:
    """Backward-compatible parsing result containing metadata and chunks."""

    doc_id: str
    metadata: dict
    chunks: List[Chunk] = field(default_factory=list)


class ParserRouter:
    """Route documents to optional backends or fallback structured parsers."""

    def __init__(self, audit_writer: ParserAuditWriter | None = None) -> None:
        self._docling_converter = DocumentConverter() if DocumentConverter else None
        self._markdown_converter = MarkItDown() if MarkItDown else None
        self.audit_writer = audit_writer or ParserAuditWriter(getattr(settings, "parser_audit_dir", "data/parser_audit"))

    @staticmethod
    def _hash_text(text: str) -> str:
        return hashlib.sha256((text or "").encode("utf-8", errors="ignore")).hexdigest()

    @staticmethod
    def _stable_chunk_id(doc_id: str, ordinal: int, text_hash: str) -> str:
        return str(uuid.uuid5(uuid.NAMESPACE_URL, f"{doc_id}:{ordinal}:{text_hash[:16]}"))

    def parse_document(
        self,
        file_path: str,
        doc_id: str | None = None,
        source_type: str = "file",
        source_url: str | None = None,
    ) -> ParsedDocument:
        """Backward-compatible wrapper around parse_document_intelligence."""
        result = self.parse_document_intelligence(
            file_path=file_path,
            doc_id=doc_id,
            source_type=source_type,
            source_url=source_url,
        )
        metadata = {
            **result.diagnostics,
            "parser": result.parser_name,
            "parser_name": result.parser_name,
            "parser_version": result.parser_version,
            "source_type": result.source_type,
            "source_url": source_url,
            "blocks_count": len(result.blocks),
            "tables_count": len(result.tables),
            "images_count": len(result.images),
            "chunks_count": len(result.chunks),
            "image_refs": [
                {
                    "url": image.source_path_or_url,
                    "alt": image.alt_text,
                    "caption": image.caption,
                    "section_path": image.section_path,
                }
                for image in result.images
            ],
            "document_intelligence": {
                "blocks": [block.model_dump() for block in result.blocks],
                "tables": [table.model_dump() for table in result.tables],
                "images": [image.model_dump() for image in result.images],
            },
        }
        return ParsedDocument(doc_id=result.doc_id, metadata=metadata, chunks=result.chunks)

    def parse_document_intelligence(
        self,
        file_path: str,
        doc_id: str | None = None,
        source_type: str = "file",
        source_url: str | None = None,
    ) -> DocumentIntelligenceResult:
        """Parse a file into blocks, tables, image refs and evidence-ready chunks."""
        path = Path(file_path)
        resolved_doc_id = doc_id or str(uuid.uuid4())
        requested = _backend_mode()
        diagnostics: dict[str, Any] = {
            "parser_backend_requested": requested,
            "parser_version": PARSER_VERSION,
            "docling_available": self._docling_converter is not None,
            "markitdown_available": self._markdown_converter is not None,
            "ocr_enabled": bool(getattr(settings, "enable_ocr", False)),
            "ocr_backend": getattr(settings, "ocr_backend", "none"),
            "warnings": [],
        }

        try:
            result = self._parse_with_backend(path, resolved_doc_id, source_type, source_url, requested, diagnostics)
        except Exception as exc:
            self.audit_writer.write_error(path.name, str(exc), diagnostics)
            raise

        scanned = detect_scanned_pdf(path, result.text)
        if path.suffix.lower() == ".pdf":
            diagnostics.update(
                {
                    "scanned_pdf_detected": scanned.is_probably_scanned,
                    "scanned_pdf_text_chars": scanned.extracted_text_chars,
                    "scanned_pdf_page_count": scanned.page_count,
                    "scanned_pdf_reason": scanned.reason,
                }
            )
            if scanned.is_probably_scanned and not getattr(settings, "enable_ocr", False):
                diagnostics["warnings"].append("PDF appears scanned; OCR is disabled")
            elif scanned.is_probably_scanned and getattr(settings, "enable_ocr", False) and ocrmypdf is None:
                diagnostics["warnings"].append("PDF appears scanned; OCR backend is unavailable")

        chunks = build_chunks_from_document(
            doc_id=resolved_doc_id,
            source_name=path.name,
            source_type=source_type,
            source_url=source_url,
            parser_name=result.parser_name,
            blocks=result.blocks,
            tables=result.tables,
            images=result.images,
            fallback_text=result.text,
        )
        result = result.model_copy(
            update={
                "chunks": chunks,
                "diagnostics": {
                    **diagnostics,
                    **result.diagnostics,
                    "parser_backend_used": result.parser_name,
                    "blocks_count": len(result.blocks),
                    "tables_count": len(result.tables),
                    "images_count": len(result.images),
                    "chunks_count": len(chunks),
                },
            }
        )
        self.audit_writer.write_parsed(result)
        return result

    def _parse_with_backend(
        self,
        path: Path,
        doc_id: str,
        source_type: str,
        source_url: str | None,
        requested: str,
        diagnostics: dict[str, Any],
    ) -> DocumentIntelligenceResult:
        ext = path.suffix.lower()
        if requested == "fallback":
            return self._parse_fallback(path, doc_id, source_type, source_url)
        if requested == "docling":
            if self._docling_converter is None:
                raise RuntimeError("PARSER_BACKEND=docling requested, but Docling is not installed")
            return self._parse_docling(path, doc_id, source_type, source_url)
        if requested == "markitdown":
            if self._markdown_converter is None:
                raise RuntimeError("PARSER_BACKEND=markitdown requested, but MarkItDown is not installed")
            return self._parse_markitdown(path, doc_id, source_type, source_url)

        if requested != "auto":
            diagnostics["warnings"].append(f"Unsupported PARSER_BACKEND={requested!r}; using auto")
        if ext == ".pdf" and self._docling_converter is not None:
            try:
                return self._parse_docling(path, doc_id, source_type, source_url)
            except Exception as exc:
                diagnostics["docling_error"] = str(exc)
                diagnostics["warnings"].append("Docling failed; fallback parser used")
        if ext in {".docx", ".pptx", ".xlsx", ".html", ".htm", ".csv", ".txt", ".md"} and self._markdown_converter is not None:
            try:
                return self._parse_markitdown(path, doc_id, source_type, source_url)
            except Exception as exc:
                diagnostics["markitdown_error"] = str(exc)
                diagnostics["warnings"].append("MarkItDown failed; fallback parser used")
        return self._parse_fallback(path, doc_id, source_type, source_url)

    def _parse_docling(self, path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
        result = self._docling_converter.convert(str(path))  # type: ignore[union-attr]
        markdown = result.document.export_to_markdown()
        parsed = _markdown_to_blocks_tables(markdown, doc_id, parser_name="docling")
        return DocumentIntelligenceResult(
            doc_id=doc_id,
            source_name=path.name,
            source_type=source_type,
            parser_name="docling",
            parser_version=PARSER_VERSION,
            text=markdown,
            blocks=parsed["blocks"],
            tables=parsed["tables"],
            images=[],
            diagnostics={"source_url": source_url},
        )

    def _parse_markitdown(self, path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
        converted = self._markdown_converter.convert(str(path))  # type: ignore[union-attr]
        markdown = converted.text_content if hasattr(converted, "text_content") else str(converted)
        parsed = _markdown_to_blocks_tables(markdown, doc_id, parser_name="markitdown")
        return DocumentIntelligenceResult(
            doc_id=doc_id,
            source_name=path.name,
            source_type=source_type,
            parser_name="markitdown",
            parser_version=PARSER_VERSION,
            text=markdown,
            blocks=parsed["blocks"],
            tables=parsed["tables"],
            images=[],
            diagnostics={"source_url": source_url},
        )

    def _parse_fallback(self, path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
        ext = path.suffix.lower()
        if ext in {".csv", ".xlsx"}:
            return _parse_spreadsheet(path, doc_id, source_type, source_url)
        if ext == ".pdf":
            return _parse_pdf(path, doc_id, source_type, source_url)
        if ext == ".docx":
            return _parse_docx(path, doc_id, source_type, source_url)
        if ext == ".pptx":
            return _parse_pptx(path, doc_id, source_type, source_url)
        if ext in {".html", ".htm"}:
            return _parse_html(path, doc_id, source_type, source_url)
        return _parse_plain_text(path, doc_id, source_type, source_url, parser_name="plain")


def table_to_row_chunks(
    table: ParsedTable,
    parser_name: str,
    source_name: str,
    source_type: str = "file",
    source_url: str | None = None,
    start_ordinal: int = 0,
) -> list[Chunk]:
    """Convert a parsed table to one evidence-ready table_row chunk per row."""
    chunks: list[Chunk] = []
    clean_headers = _clean_headers(table.headers)
    columns_text = " | ".join(clean_headers)
    for row_index, row in enumerate(table.rows):
        parts = []
        for column, value in zip(clean_headers, row):
            value_text = str(value or "").strip()
            if value_text and value_text.lower() != "nan":
                parts.append(f"{column}: {value_text}")
        if not parts:
            continue
        text = "\n".join(
            part
            for part in [
                f"Table: {table.caption or table.table_id}",
                f"Table columns: {columns_text}",
                " | ".join(parts),
            ]
            if part
        )
        original_text_hash = _hash_text(text)
        normalized_text = normalize_dirty_scientific_text(text)
        normalization_applied = normalized_text != text
        text = normalized_text
        text_hash = _hash_text(text)
        ordinal = start_ordinal + len(chunks)
        chunks.append(
            Chunk(
                chunk_id=_stable_chunk_id(table.doc_id, table.table_id, row_index, text_hash),
                doc_id=table.doc_id,
                text=text,
                page_start=table.page_start or 1,
                page_end=table.page_end or table.page_start or 1,
                section_path=table.section_path or "/",
                ordinal=ordinal,
                char_start=None,
                char_end=None,
                token_count=len(re.findall(r"\S+", text)),
                text_hash=text_hash,
                preview=_preview(text),
                embedding_version=settings.embedding_model,
                metadata={
                    "parser": parser_name,
                    "parser_name": parser_name,
                    "source_name": source_name,
                    "filename": source_name,
                    "source_type": source_type,
                    "source_url": source_url,
                    "chunk_kind": "table_row",
                    "block_type": "table_row",
                    "table_id": table.table_id,
                    "row_id": str(row_index),
                    "table_columns": columns_text,
                    "table_caption": table.caption,
                    "source_block_id": table.source_block_id,
                    "parser_version": PARSER_VERSION,
                    "text_normalization_applied": normalization_applied,
                    "original_text_hash": original_text_hash if normalization_applied else None,
                },
            )
        )
    return chunks


def build_chunks_from_document(
    doc_id: str,
    source_name: str,
    source_type: str,
    source_url: str | None,
    parser_name: str,
    blocks: list[ParsedBlock],
    tables: list[ParsedTable],
    images: list[ParsedImageRef],
    fallback_text: str,
) -> list[Chunk]:
    """Build structured chunks from blocks/tables/images with stable provenance."""
    chunks: list[Chunk] = []
    for table in tables:
        chunks.extend(
            table_to_row_chunks(
                table,
                parser_name=parser_name,
                source_name=source_name,
                source_type=source_type,
                source_url=source_url,
                start_ordinal=len(chunks),
            )
        )

    table_block_ids = {table.source_block_id for table in tables if table.source_block_id}
    for block in blocks:
        if block.block_id in table_block_ids or block.block_type == "table":
            continue
        chunks.extend(
            _block_to_chunks(
                block=block,
                parser_name=parser_name,
                source_name=source_name,
                source_type=source_type,
                source_url=source_url,
                start_ordinal=len(chunks),
            )
        )

    for image in images:
        text = f"Image: url: {image.source_path_or_url or ''} | alt: {image.alt_text or ''} | caption: {image.caption or ''} | section_path: {image.section_path}"
        if not text.strip():
            continue
        text_hash = _hash_text(text)
        chunks.append(
            Chunk(
                chunk_id=_stable_chunk_id(doc_id, image.image_id, 0, text_hash),
                doc_id=doc_id,
                text=text,
                page_start=image.page or 1,
                page_end=image.page or 1,
                section_path=image.section_path or "/",
                ordinal=len(chunks),
                token_count=len(re.findall(r"\S+", text)),
                text_hash=text_hash,
                preview=_preview(text),
                embedding_version=settings.embedding_model,
                metadata={
                    "parser": parser_name,
                    "parser_name": parser_name,
                    "source_name": source_name,
                    "filename": source_name,
                    "source_type": source_type,
                    "source_url": source_url,
                    "chunk_kind": "image_ref",
                    "block_type": "figure",
                    "image_id": image.image_id,
                    "image_refs": [
                        {
                            "url": image.source_path_or_url,
                            "alt": image.alt_text,
                            "caption": image.caption,
                            "section_path": image.section_path,
                        }
                    ],
                    "parser_version": PARSER_VERSION,
                },
            )
        )

    if not chunks and fallback_text.strip():
        fallback_block = ParsedBlock(
            block_id=_stable_id("block", doc_id, "fallback", _hash_text(fallback_text)),
            doc_id=doc_id,
            block_type="paragraph",
            text=fallback_text.strip(),
            page_start=1,
            page_end=1,
            section_path="/",
            ordinal=0,
            metadata={},
        )
        chunks.extend(
            _block_to_chunks(
                block=fallback_block,
                parser_name=parser_name,
                source_name=source_name,
                source_type=source_type,
                source_url=source_url,
                start_ordinal=0,
            )
        )

    return [chunk for idx, chunk in enumerate(chunks) if _renumber_chunk(chunk, idx)]


def detect_scanned_pdf(file_path: str | Path, extracted_text: str) -> ScannedPdfDetection:
    """Detect likely scanned PDFs without making OCR mandatory."""
    path = Path(file_path)
    if path.suffix.lower() != ".pdf":
        return ScannedPdfDetection(is_probably_scanned=False, extracted_text_chars=len(extracted_text or ""), page_count=None, reason="not_pdf")
    page_count: int | None = None
    try:
        from pypdf import PdfReader

        page_count = len(PdfReader(str(path)).pages)
    except Exception:
        page_count = None
    chars = len((extracted_text or "").strip())
    threshold = int(getattr(settings, "scanned_pdf_min_text_chars", 50))
    is_scanned = bool(page_count and page_count > 0 and chars < threshold)
    reason = f"text_chars_below_{threshold}" if is_scanned else "enough_text_or_unknown_pages"
    return ScannedPdfDetection(is_probably_scanned=is_scanned, extracted_text_chars=chars, page_count=page_count, reason=reason)


def _parse_spreadsheet(path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
    import pandas as pd

    tables: list[ParsedTable] = []
    blocks: list[ParsedBlock] = []
    if path.suffix.lower() == ".xlsx":
        loaded = pd.read_excel(path, sheet_name=None)
        frames = list(loaded.items()) if isinstance(loaded, dict) else [("Sheet1", loaded)]
        parser_name = "pandas"
    else:
        try:
            frame = pd.read_csv(path, sep=None, engine="python", encoding="utf-8-sig")
        except Exception:
            frame = pd.read_csv(path, encoding="utf-8-sig")
        frames = [(path.stem, frame)]
        parser_name = "pandas"
    for idx, (sheet_name, frame) in enumerate(frames, start=1):
        headers = [str(column) for column in frame.columns]
        rows = [[_cell(value) for value in row.tolist()] for _, row in frame.iterrows()]
        table_id = _stable_id("table", doc_id, sheet_name, idx)
        block_id = _stable_id("block", doc_id, "table", sheet_name, idx)
        blocks.append(
            ParsedBlock(
                block_id=block_id,
                doc_id=doc_id,
                block_type="table",
                text="\n".join(_format_table_rows(headers, rows, table_name=str(sheet_name))),
                page_start=1,
                page_end=1,
                section_path="/",
                ordinal=len(blocks),
                metadata={"sheet_name": str(sheet_name)},
            )
        )
        tables.append(
            ParsedTable(
                table_id=table_id,
                doc_id=doc_id,
                page_start=1,
                page_end=1,
                section_path="/",
                caption=str(sheet_name),
                headers=headers,
                rows=rows,
                source_block_id=block_id,
                metadata={"sheet_name": str(sheet_name)},
            )
        )
    text = "\n\n".join(block.text for block in blocks)
    return DocumentIntelligenceResult(
        doc_id=doc_id,
        source_name=path.name,
        source_type=source_type,
        parser_name=parser_name,
        parser_version=PARSER_VERSION,
        text=text,
        blocks=blocks,
        tables=tables,
        diagnostics={"source_url": source_url},
    )


def _parse_pdf(path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    blocks: list[ParsedBlock] = []
    texts: list[str] = []
    for page_idx, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        texts.append(page_text)
        for paragraph in _paragraphs(page_text):
            blocks.append(
                ParsedBlock(
                    block_id=_stable_id("block", doc_id, "pdf", page_idx, len(blocks), _hash_text(paragraph)),
                    doc_id=doc_id,
                    block_type="paragraph",
                    text=paragraph,
                    page_start=page_idx,
                    page_end=page_idx,
                    section_path="/",
                    ordinal=len(blocks),
                    metadata={"page": page_idx},
                )
            )
    text = "\n\n".join(texts).strip()
    return DocumentIntelligenceResult(
        doc_id=doc_id,
        source_name=path.name,
        source_type=source_type,
        parser_name="pypdf",
        parser_version=PARSER_VERSION,
        text=text,
        blocks=blocks,
        diagnostics={"source_url": source_url},
    )


def _parse_docx(path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
    import docx

    document = docx.Document(str(path))
    blocks: list[ParsedBlock] = []
    tables: list[ParsedTable] = []
    section_path = "/"
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = str(getattr(paragraph.style, "name", "") or "")
        block_type = "section_heading" if style_name.lower().startswith("heading") or "заголовок" in style_name.lower() else "paragraph"
        if block_type == "section_heading":
            section_path = _join_section(section_path, text)
        blocks.append(
            ParsedBlock(
                block_id=_stable_id("block", doc_id, "docx", len(blocks), _hash_text(text)),
                doc_id=doc_id,
                block_type=block_type,
                text=text,
                page_start=1,
                page_end=1,
                section_path=section_path,
                ordinal=len(blocks),
                metadata={"style": style_name},
            )
        )
    for idx, table in enumerate(document.tables, start=1):
        raw_rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not raw_rows:
            continue
        headers = raw_rows[0]
        rows = raw_rows[1:] if len(raw_rows) > 1 else []
        table_id = _stable_id("table", doc_id, "docx", idx)
        block_id = _stable_id("block", doc_id, "docx-table", idx)
        caption = f"DOCX table {idx}"
        blocks.append(
            ParsedBlock(
                block_id=block_id,
                doc_id=doc_id,
                block_type="table",
                text="\n".join(_format_table_rows(headers, rows, table_name=caption)),
                page_start=1,
                page_end=1,
                section_path=section_path,
                ordinal=len(blocks),
                metadata={},
            )
        )
        tables.append(
            ParsedTable(
                table_id=table_id,
                doc_id=doc_id,
                page_start=1,
                page_end=1,
                section_path=section_path,
                caption=caption,
                headers=headers,
                rows=rows,
                source_block_id=block_id,
            )
        )
    text = "\n\n".join(block.text for block in blocks)
    return DocumentIntelligenceResult(
        doc_id=doc_id,
        source_name=path.name,
        source_type=source_type,
        parser_name="python-docx",
        parser_version=PARSER_VERSION,
        text=text,
        blocks=blocks,
        tables=tables,
        diagnostics={"source_url": source_url},
    )


def _parse_pptx(path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
    from pptx import Presentation  # type: ignore

    presentation = Presentation(str(path))
    blocks: list[ParsedBlock] = []
    tables: list[ParsedTable] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        section_path = f"/Slide {slide_idx}"
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    blocks.append(
                        ParsedBlock(
                            block_id=_stable_id("block", doc_id, "pptx", slide_idx, shape_idx, _hash_text(text)),
                            doc_id=doc_id,
                            block_type="paragraph",
                            text=f"Slide {slide_idx}: {text}",
                            page_start=slide_idx,
                            page_end=slide_idx,
                            section_path=section_path,
                            ordinal=len(blocks),
                            metadata={"slide": slide_idx},
                        )
                    )
            if getattr(shape, "has_table", False):
                raw_rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                if not raw_rows:
                    continue
                headers = raw_rows[0]
                rows = raw_rows[1:] if len(raw_rows) > 1 else []
                table_id = _stable_id("table", doc_id, "pptx", slide_idx, shape_idx)
                block_id = _stable_id("block", doc_id, "pptx-table", slide_idx, shape_idx)
                caption = f"PPTX slide {slide_idx} table"
                blocks.append(
                    ParsedBlock(
                        block_id=block_id,
                        doc_id=doc_id,
                        block_type="table",
                        text="\n".join(_format_table_rows(headers, rows, table_name=caption)),
                        page_start=slide_idx,
                        page_end=slide_idx,
                        section_path=section_path,
                        ordinal=len(blocks),
                        metadata={"slide": slide_idx},
                    )
                )
                tables.append(
                    ParsedTable(
                        table_id=table_id,
                        doc_id=doc_id,
                        page_start=slide_idx,
                        page_end=slide_idx,
                        section_path=section_path,
                        caption=caption,
                        headers=headers,
                        rows=rows,
                        source_block_id=block_id,
                        metadata={"slide": slide_idx},
                    )
                )
    text = "\n\n".join(block.text for block in blocks)
    return DocumentIntelligenceResult(
        doc_id=doc_id,
        source_name=path.name,
        source_type=source_type,
        parser_name="python-pptx",
        parser_version=PARSER_VERSION,
        text=text,
        blocks=blocks,
        tables=tables,
        diagnostics={"source_url": source_url},
    )


def _parse_html(path: Path, doc_id: str, source_type: str, source_url: str | None) -> DocumentIntelligenceResult:
    from bs4 import BeautifulSoup  # type: ignore

    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    for node in soup(["script", "style", "nav", "footer"]):
        node.extract()
    blocks: list[ParsedBlock] = []
    tables: list[ParsedTable] = []
    images: list[ParsedImageRef] = []
    section_path = "/"
    if soup.title and soup.title.get_text(" ", strip=True):
        title = soup.title.get_text(" ", strip=True)
        blocks.append(
            ParsedBlock(
                block_id=_stable_id("block", doc_id, "title", _hash_text(title)),
                doc_id=doc_id,
                block_type="title",
                text=title,
                page_start=1,
                page_end=1,
                section_path="/",
                ordinal=len(blocks),
            )
        )
    for node in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "table", "figure", "img", "figcaption"]):
        name = node.name.lower()
        if name in {"h1", "h2", "h3", "h4"}:
            text = node.get_text(" ", strip=True)
            if not text:
                continue
            section_path = _join_section("/", text) if name == "h1" else _join_section(section_path, text)
            blocks.append(_block(doc_id, "section_heading", text, len(blocks), section_path))
        elif name in {"p", "li"}:
            text = node.get_text(" ", strip=True)
            if text:
                blocks.append(_block(doc_id, "list_item" if name == "li" else "paragraph", text, len(blocks), section_path))
        elif name == "figcaption":
            text = node.get_text(" ", strip=True)
            if text:
                blocks.append(_block(doc_id, "caption", text, len(blocks), section_path))
        elif name == "img":
            images.append(_image_from_html(node, doc_id, len(images), section_path))
        elif name == "figure":
            img = node.find("img")
            if img:
                images.append(_image_from_html(img, doc_id, len(images), section_path, figure=node))
        elif name == "table":
            rows = []
            for tr in node.find_all("tr"):
                cells = tr.find_all(["th", "td"])
                row = [cell.get_text(" ", strip=True) for cell in cells]
                if any(row):
                    rows.append(row)
            if rows:
                headers = rows[0]
                body = rows[1:] if len(rows) > 1 else []
                caption_node = node.find("caption")
                caption = caption_node.get_text(" ", strip=True) if caption_node else f"HTML table {len(tables) + 1}"
                block_id = _stable_id("block", doc_id, "html-table", len(tables), _hash_text(caption))
                table_id = _stable_id("table", doc_id, "html", len(tables), _hash_text(caption))
                blocks.append(
                    ParsedBlock(
                        block_id=block_id,
                        doc_id=doc_id,
                        block_type="table",
                        text="\n".join(_format_table_rows(headers, body, table_name=caption)),
                        page_start=1,
                        page_end=1,
                        section_path=section_path,
                        ordinal=len(blocks),
                    )
                )
                tables.append(
                    ParsedTable(
                        table_id=table_id,
                        doc_id=doc_id,
                        page_start=1,
                        page_end=1,
                        section_path=section_path,
                        caption=caption,
                        headers=headers,
                        rows=body,
                        source_block_id=block_id,
                    )
                )
    text = "\n\n".join(block.text for block in blocks)
    return DocumentIntelligenceResult(
        doc_id=doc_id,
        source_name=path.name,
        source_type=source_type,
        parser_name="html",
        parser_version=PARSER_VERSION,
        text=text,
        blocks=blocks,
        tables=tables,
        images=_dedupe_images(images),
        diagnostics={"source_url": source_url, "title": soup.title.get_text(" ", strip=True) if soup.title else None},
    )


def _parse_plain_text(path: Path, doc_id: str, source_type: str, source_url: str | None, parser_name: str) -> DocumentIntelligenceResult:
    text = path.read_text(encoding="utf-8", errors="ignore")
    parsed = _markdown_to_blocks_tables(text, doc_id, parser_name=parser_name)
    return DocumentIntelligenceResult(
        doc_id=doc_id,
        source_name=path.name,
        source_type=source_type,
        parser_name=parser_name,
        parser_version=PARSER_VERSION,
        text=text,
        blocks=parsed["blocks"],
        tables=parsed["tables"],
        diagnostics={"source_url": source_url},
    )


def _markdown_to_blocks_tables(text: str, doc_id: str, parser_name: str) -> dict[str, list[Any]]:
    blocks: list[ParsedBlock] = []
    tables: list[ParsedTable] = []
    section_path = "/"
    lines = text.splitlines()
    paragraph_buffer: list[str] = []

    def flush_paragraph() -> None:
        nonlocal paragraph_buffer
        paragraph = " ".join(part.strip() for part in paragraph_buffer if part.strip()).strip()
        paragraph_buffer = []
        if paragraph:
            blocks.append(_block(doc_id, "paragraph", paragraph, len(blocks), section_path, metadata={"parser_hint": parser_name}))

    idx = 0
    while idx < len(lines):
        line = lines[idx].strip()
        if not line:
            flush_paragraph()
            idx += 1
            continue
        if line.startswith("#"):
            flush_paragraph()
            heading = line.lstrip("#").strip()
            section_path = _join_section("/", heading) if line.startswith("# ") else _join_section(section_path, heading)
            blocks.append(_block(doc_id, "section_heading", heading, len(blocks), section_path))
            idx += 1
            continue
        if "|" in line and idx + 1 < len(lines) and re.match(r"^\s*\|?\s*:?-{2,}", lines[idx + 1]):
            flush_paragraph()
            table_lines = [line]
            idx += 2
            while idx < len(lines) and "|" in lines[idx]:
                table_lines.append(lines[idx].strip())
                idx += 1
            headers, rows = _parse_markdown_table(table_lines)
            if headers:
                table_id = _stable_id("table", doc_id, "md", len(tables), _hash_text("\n".join(table_lines)))
                block_id = _stable_id("block", doc_id, "md-table", len(tables), _hash_text("\n".join(table_lines)))
                blocks.append(
                    ParsedBlock(
                        block_id=block_id,
                        doc_id=doc_id,
                        block_type="table",
                        text="\n".join(_format_table_rows(headers, rows, table_name=f"Markdown table {len(tables) + 1}")),
                        page_start=1,
                        page_end=1,
                        section_path=section_path,
                        ordinal=len(blocks),
                    )
                )
                tables.append(
                    ParsedTable(
                        table_id=table_id,
                        doc_id=doc_id,
                        page_start=1,
                        page_end=1,
                        section_path=section_path,
                        caption=f"Markdown table {len(tables) + 1}",
                        headers=headers,
                        rows=rows,
                        source_block_id=block_id,
                    )
                )
            continue
        paragraph_buffer.append(line)
        idx += 1
    flush_paragraph()
    return {"blocks": blocks, "tables": tables}


def _block(doc_id: str, block_type: str, text: str, ordinal: int, section_path: str, metadata: dict[str, Any] | None = None) -> ParsedBlock:
    return ParsedBlock(
        block_id=_stable_id("block", doc_id, block_type, ordinal, _hash_text(text)),
        doc_id=doc_id,
        block_type=block_type,  # type: ignore[arg-type]
        text=text,
        page_start=1,
        page_end=1,
        section_path=section_path or "/",
        ordinal=ordinal,
        metadata=metadata or {},
    )


def _block_to_chunks(
    block: ParsedBlock,
    parser_name: str,
    source_name: str,
    source_type: str,
    source_url: str | None,
    start_ordinal: int,
) -> list[Chunk]:
    text = (block.text or "").strip()
    if not text:
        return []
    words = list(re.finditer(r"\S+", text))
    if not words:
        return []
    size = max(50, int(settings.chunk_size))
    overlap = max(0, min(int(settings.chunk_overlap), size // 2))
    step = max(1, size - overlap)
    if len(words) <= size:
        return [
            _make_chunk_from_text(
                doc_id=block.doc_id,
                text=text,
                parser_name=parser_name,
                source_name=source_name,
                source_type=source_type,
                source_url=source_url,
                ordinal=start_ordinal,
                page_start=block.page_start or 1,
                page_end=block.page_end or block.page_start or 1,
                section_path=block.section_path,
                chunk_kind=block.block_type if block.block_type in {"paragraph", "section_heading", "caption", "list_item"} else "paragraph",
                block_type=block.block_type,
                block_id=block.block_id,
                char_start=0,
                char_end=len(text),
                extra_metadata=block.metadata,
            )
        ]
    chunks: list[Chunk] = []
    cursor = 0
    while cursor < len(words):
        window = words[cursor : cursor + size]
        chunk_text = text[window[0].start() : window[-1].end()].strip()
        chunks.append(
            _make_chunk_from_text(
                doc_id=block.doc_id,
                text=chunk_text,
                parser_name=parser_name,
                source_name=source_name,
                source_type=source_type,
                source_url=source_url,
                ordinal=start_ordinal + len(chunks),
                page_start=block.page_start or 1,
                page_end=block.page_end or block.page_start or 1,
                section_path=block.section_path,
                chunk_kind="text_window",
                block_type=block.block_type,
                block_id=block.block_id,
                char_start=window[0].start(),
                char_end=window[-1].end(),
                extra_metadata=block.metadata,
            )
        )
        cursor += step
    return chunks


def _make_chunk_from_text(
    doc_id: str,
    text: str,
    parser_name: str,
    source_name: str,
    source_type: str,
    source_url: str | None,
    ordinal: int,
    page_start: int,
    page_end: int,
    section_path: str,
    chunk_kind: str,
    block_type: str,
    block_id: str,
    char_start: int | None,
    char_end: int | None,
    extra_metadata: dict[str, Any] | None = None,
) -> Chunk:
    original_text_hash = _hash_text(text)
    normalized_text = normalize_dirty_scientific_text(text)
    normalization_applied = normalized_text != text
    text = normalized_text
    text_hash = _hash_text(text)
    return Chunk(
        chunk_id=_stable_chunk_id(doc_id, block_id, ordinal, text_hash),
        doc_id=doc_id,
        text=text,
        page_start=page_start,
        page_end=page_end,
        section_path=section_path or "/",
        ordinal=ordinal,
        char_start=char_start,
        char_end=char_end,
        token_count=len(re.findall(r"\S+", text)),
        text_hash=text_hash,
        preview=_preview(text),
        embedding_version=settings.embedding_model,
        metadata={
            **(extra_metadata or {}),
            "parser": parser_name,
            "parser_name": parser_name,
            "source_name": source_name,
            "filename": source_name,
            "source_type": source_type,
            "source_url": source_url,
            "chunk_kind": chunk_kind,
            "block_type": block_type,
            "block_id": block_id,
            "parser_version": PARSER_VERSION,
            "text_normalization_applied": normalization_applied,
            "original_text_hash": original_text_hash if normalization_applied else None,
        },
    )


def _renumber_chunk(chunk: Chunk, ordinal: int) -> bool:
    if not chunk.text.strip():
        return False
    chunk.ordinal = ordinal
    return True


def _parse_markdown_table(lines: list[str]) -> tuple[list[str], list[list[str]]]:
    if not lines:
        return [], []
    headers = [cell.strip() for cell in lines[0].strip("|").split("|")]
    rows = [[cell.strip() for cell in line.strip("|").split("|")] for line in lines[1:]]
    return headers, rows


def _format_table_rows(headers: List[str], rows: List[List[Any]], table_name: str | None = None) -> List[str]:
    clean_headers = _clean_headers(headers)
    lines: list[str] = []
    if table_name:
        lines.append(f"Table: {table_name}")
    if clean_headers:
        lines.append("Table columns: " + " | ".join(clean_headers))
    for row in rows:
        parts = []
        for col_name, value in zip(clean_headers, row):
            value_text = str(value or "").strip()
            if value_text and value_text.lower() != "nan":
                parts.append(f"{col_name}: {value_text}")
        if parts:
            lines.append(" | ".join(parts))
    return lines


def _clean_headers(headers: Iterable[Any]) -> list[str]:
    return [str(header).strip() if str(header).strip() else f"column_{idx + 1}" for idx, header in enumerate(headers)]


def _cell(value: Any) -> str:
    try:
        import pandas as pd

        if pd.isna(value):
            return ""
    except Exception:
        pass
    return str(value).strip()


def _paragraphs(text: str) -> list[str]:
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n+", text or "") if part.strip()]
    if paragraphs:
        return paragraphs
    return [line.strip() for line in (text or "").splitlines() if line.strip()]


def _image_from_html(node, doc_id: str, idx: int, section_path: str, figure=None) -> ParsedImageRef:
    caption = ""
    parent = figure or node.find_parent("figure")
    if parent:
        caption_node = parent.find("figcaption")
        caption = caption_node.get_text(" ", strip=True) if caption_node else ""
    src = node.get("src") or ""
    alt = node.get("alt") or ""
    return ParsedImageRef(
        image_id=_stable_id("image", doc_id, idx, src, alt, caption),
        doc_id=doc_id,
        page=1,
        section_path=section_path or "/",
        alt_text=alt,
        caption=caption,
        source_path_or_url=src,
    )


def _dedupe_images(images: list[ParsedImageRef]) -> list[ParsedImageRef]:
    seen = set()
    result = []
    for image in images:
        key = (image.source_path_or_url, image.alt_text, image.caption, image.section_path)
        if key in seen:
            continue
        seen.add(key)
        result.append(image)
    return result


def _join_section(current: str, heading: str) -> str:
    cleaned = re.sub(r"\s+", " ", heading or "").strip(" /")
    if not cleaned:
        return current or "/"
    base = "/" if not current or current == "/" else current.rstrip("/")
    return f"{base}/{cleaned}" if base != "/" else f"/{cleaned}"


def _preview(text: str) -> str:
    return text[:200] + ("..." if len(text) > 200 else "")


def _hash_text(text: str) -> str:
    return hashlib.sha256((text or "").encode("utf-8", errors="ignore")).hexdigest()


def _stable_id(prefix: str, *parts: object) -> str:
    raw = "|".join("" if part is None else str(part) for part in parts)
    return f"{prefix}_{hashlib.sha1(raw.encode('utf-8')).hexdigest()[:24]}"


def _stable_chunk_id(doc_id: str, anchor: str, ordinal: int, text_hash: str) -> str:
    return str(uuid.uuid5(uuid.NAMESPACE_URL, f"{doc_id}:{anchor}:{ordinal}:{text_hash[:16]}"))


def _backend_mode() -> str:
    mode = str(getattr(settings, "parser_backend", "auto") or "auto").lower()
    return mode if mode in {"auto", "docling", "markitdown", "fallback"} else "auto"
