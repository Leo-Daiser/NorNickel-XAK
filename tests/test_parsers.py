from __future__ import annotations

from pathlib import Path

from app.ingestion.parser_router import ParserRouter


def _joined_text(chunks) -> str:
    return "\n".join(chunk.text for chunk in chunks)


def test_parser_csv_preserves_headers(tmp_path: Path) -> None:
    csv_path = tmp_path / "experiments.csv"
    csv_path.write_text(
        "material,process_regime,property,value,unit,equipment,laboratory\n"
        "ВТ6,отжиг,прочность,980,MPa,Вакуумная печь SNOL-75,Лаборатория легких сплавов\n",
        encoding="utf-8",
    )

    parsed = ParserRouter().parse_document(str(csv_path), doc_id="csv-doc")
    text = _joined_text(parsed.chunks)

    assert parsed.chunks
    assert "Table columns:" in text
    assert "material: ВТ6" in text
    assert "process_regime: отжиг" in text


def test_parser_xlsx_preserves_headers(tmp_path: Path) -> None:
    from openpyxl import Workbook

    xlsx_path = tmp_path / "experiments.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "experiments"
    ws.append(["material", "process_regime", "property", "value", "unit"])
    ws.append(["сталь 12Х18Н10Т", "закалка", "твёрдость", 210, "HV"])
    wb.save(xlsx_path)

    parsed = ParserRouter().parse_document(str(xlsx_path), doc_id="xlsx-doc")
    text = _joined_text(parsed.chunks)

    assert parsed.chunks
    assert "Table: experiments" in text
    assert "material: сталь 12Х18Н10Т" in text
    assert "property: твёрдость" in text


def test_parser_docx_extracts_tables(tmp_path: Path) -> None:
    import docx

    docx_path = tmp_path / "article.docx"
    doc = docx.Document()
    doc.add_paragraph("Experiment: DOCX-VT6")
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "material"
    table.cell(0, 1).text = "process_regime"
    table.cell(0, 2).text = "property"
    table.cell(1, 0).text = "VT6"
    table.cell(1, 1).text = "aging 550 C 4 h"
    table.cell(1, 2).text = "hardness 360 HV increased"
    doc.save(docx_path)

    parsed = ParserRouter().parse_document(str(docx_path), doc_id="docx-doc")
    text = _joined_text(parsed.chunks)

    assert "Table: DOCX table 1" in text
    assert "material: VT6" in text
    assert "process_regime: aging 550 C 4 h" in text


def test_parser_pptx_extracts_text_and_tables(tmp_path: Path) -> None:
    from pptx import Presentation

    pptx_path = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Experiment: PPTX-STEEL"
    table_shape = slide.shapes.add_table(2, 3, 100000, 1000000, 8000000, 1000000)
    table = table_shape.table
    table.cell(0, 0).text = "material"
    table.cell(0, 1).text = "process_regime"
    table.cell(0, 2).text = "property"
    table.cell(1, 0).text = "сталь 12Х18Н10Т"
    table.cell(1, 1).text = "закалка 1050 C 1 h"
    table.cell(1, 2).text = "твёрдость 210 HV increased"
    prs.save(pptx_path)

    parsed = ParserRouter().parse_document(str(pptx_path), doc_id="pptx-doc")
    text = _joined_text(parsed.chunks)

    assert "Experiment: PPTX-STEEL" in text
    assert "Table: PPTX slide 1 table" in text
    assert "material: сталь 12Х18Н10Т" in text


def test_parser_html_extracts_tables(tmp_path: Path) -> None:
    html_path = tmp_path / "report.html"
    html_path.write_text(
        "<html><body><table>"
        "<tr><th>material</th><th>process_regime</th><th>property</th></tr>"
        "<tr><td>7075</td><td>старение 160 C 8 h</td><td>коррозионная стойкость</td></tr>"
        "</table></body></html>",
        encoding="utf-8",
    )

    parsed = ParserRouter().parse_document(str(html_path), doc_id="html-doc")
    text = _joined_text(parsed.chunks)

    assert "Table: HTML table 1" in text
    assert "material: 7075" in text
    assert "коррозионная стойкость" in text
